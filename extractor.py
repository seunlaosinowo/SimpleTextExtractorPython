import sys

x = sys.argv[1]
y = sys.argv[2]

csvfile="/var/www/splattng/uploads/"+str(y)
pptfile="/var/www/splattng/uploads/"+str(x)


from pptx import Presentation

prs = Presentation(pptfile)
f = open(csvfile, "w", encoding='utf-16')

for eachslide in prs.slides:
    for item in eachslide.placeholders:
            if item.placeholder_format.idx==0:
                f.write(item.text+",")
                
                
            if item.placeholder_format.idx==1:
                cont=item.text.rstrip("\n")
                cont=cont.replace('\n','#')
                f.write('\"'+cont+'\"'+"\n")
              
f.close()
#print (pptfile)


